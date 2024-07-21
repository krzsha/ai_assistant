from flask import Flask, request, render_template
from PyPDF2 import PdfReader
from pptx import Presentation
from transformers import pipeline
from docx import Document
from tqdm import tqdm
app = Flask(__name__)

# Functions to parse different document types
def parse_pdf(file):
    print('read')
    reader = PdfReader(file)
    text = ''
    for page in range(len(reader.pages)):
        text += reader.pages[page].extract_text()
    return text

def parse_ppt(file):
    prs = Presentation(file)
    text = ''
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text
    return text

def parse_word(file):
    doc = Document(file)
    text = ''
    for paragraph in doc.paragraphs:
        text += paragraph.text
    return text

# Summarization pipeline
summarizer = pipeline("summarization", model="facebook/bart-large-cnn")

def summarize_text(text):
    # Split text into smaller chunks if it's too long
    max_chunk_size = 1024
    text_chunks = [text[i:i+max_chunk_size] for i in range(0, len(text), max_chunk_size)]
    
    summary = ''
    for chunk in tqdm(text_chunks):
        print("i got it")
        summary += summarizer(chunk, max_length=150, min_length=30, do_sample=False)[0]['summary_text'] + ' '
    return summary.strip()

# Question-Answering pipeline
qa_pipeline = pipeline("question-answering", model="distilbert-base-cased-distilled-squad")

def answer_question(context, question):
    answer = qa_pipeline(question=question, context=context)
    return answer['answer']

@app.route('/')
def index():
    return render_template('index.html')

@app.route('/upload', methods=['POST'])
def upload():
    uploaded_files = request.files.getlist('files')
    docs = []
    for file in uploaded_files:
        if file.filename.endswith('.pdf'):
            docs.append(parse_pdf(file))
        elif file.filename.endswith('.pptx'):
            docs.append(parse_ppt(file))
        elif file.filename.endswith('.docx'):
            docs.append(parse_word(file))
    print('goot te')
    summaries = [summarize_text(doc) for doc in docs]
    context = " ".join(docs)
    
    user_query = request.form.get('query')
    answer = ''
    if user_query:
        answer = answer_question(context, user_query)

    return render_template('results.html', summaries=summaries, answer=answer)

if __name__ == '__main__':
    app.run(debug=True)

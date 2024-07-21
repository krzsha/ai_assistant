import streamlit as st
from PyPDF2 import PdfReader
from pptx import Presentation
from transformers import pipeline
from docx import Document
import torch
from concurrent.futures import ThreadPoolExecutor

# Functions to parse different document types
def parse_pdf(file):
    reader = PdfReader(file)
    text = ''
    for page in range(len(reader.pages)):
        text += reader.pages[page].extract_text()
    return text

def parse_ppt(file):
    prs = Presentation(file)
    text = ''
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text
    return text

def parse_word(file):
    doc = Document(file)
    text = ''
    for paragraph in doc.paragraphs:
        text += paragraph.text
    return text

# Check if GPU is available
device = 'cuda' if torch.cuda.is_available() else 'cpu'

# Summarization pipeline
summarizer = pipeline("summarization", model="facebook/bart-large-cnn", device=0 if device == 'cuda' else -1)

def summarize_text_chunk(chunk):
    return summarizer(chunk, max_length=150, min_length=30, do_sample=False)[0]['summary_text']

def summarize_text(text):
    # Split text into smaller chunks if it's too long
    max_chunk_size = 512
    text_chunks = [text[i:i+max_chunk_size] for i in range(0, len(text), max_chunk_size)]
    
    summary = ''
    with ThreadPoolExecutor() as executor:
        futures = [executor.submit(summarize_text_chunk, chunk) for chunk in text_chunks]
        for future in futures:
            summary += future.result() + ' '
    
    return summary.strip()

# Question-Answering pipeline
qa_pipeline = pipeline("question-answering", model="distilbert-base-cased-distilled-squad", device=0 if device == 'cuda' else -1)

def answer_question(context, question):
    answer = qa_pipeline(question=question, context=context)
    return answer['answer']

# Streamlit interface
st.title('AI Assistant')

uploaded_files = st.file_uploader("Choose files", accept_multiple_files=True)
if uploaded_files:
    docs = []
    for file in uploaded_files:
        if file.name.endswith('.pdf'):
            docs.append(parse_pdf(file))
        elif file.name.endswith('.pptx'):
            docs.append(parse_ppt(file))
        elif file.name.endswith('.docx'):
            docs.append(parse_word(file))

    st.write('Documents uploaded successfully!')

    for doc in docs:
        summary = summarize_text(doc)
        st.write('Summary:', summary)

    user_query = st.text_input("Ask a question:")
    if user_query:
        context = " ".join(docs)
        answer = answer_question(context, user_query)
        st.write('Answer:', answer)
